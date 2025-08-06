#!/usr/bin/env python3
"""
Robust PDF to PowerPoint Converter

A production-ready version of the PDF to PowerPoint converter with comprehensive
error handling, input validation, logging, and configuration options.
"""

import argparse
import logging
import os
import sys
from pathlib import Path
from typing import Optional, Tuple
import tempfile
import shutil

from PIL import Image
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO


class PDFToPowerPointConverter:
    """Robust PDF to PowerPoint converter with comprehensive error handling."""
    
    def __init__(self, dpi: int = 300, image_format: str = 'png', 
                 thread_count: int = 2, progress_interval: int = 10):
        """
        Initialize the converter with configuration options.
        
        Args:
            dpi: Resolution for PDF to image conversion (default: 300)
            image_format: Image format for slides ('png', 'jpeg', 'tiff')
            thread_count: Number of threads for PDF conversion
            progress_interval: How often to show progress (every N slides)
        """
        self.dpi = dpi
        self.image_format = image_format.lower()
        self.thread_count = thread_count
        self.progress_interval = progress_interval
        
        if self.image_format not in ['png', 'jpeg', 'tiff']:
            raise ValueError(f"Unsupported image format: {image_format}")
        
        self.logger = logging.getLogger(__name__)
        
    def validate_input_file(self, pdf_path: str) -> Path:
        """
        Validate the input PDF file.
        
        Args:
            pdf_path: Path to the PDF file
            
        Returns:
            Path object for the validated file
            
        Raises:
            FileNotFoundError: If file doesn't exist
            ValueError: If file is not a PDF or is invalid
        """
        pdf_file = Path(pdf_path)
        
        if not pdf_file.exists():
            raise FileNotFoundError(f"PDF file not found: {pdf_path}")
        
        if not pdf_file.is_file():
            raise ValueError(f"Path is not a file: {pdf_path}")
        
        if pdf_file.suffix.lower() != '.pdf':
            raise ValueError(f"File is not a PDF: {pdf_path}")
        
        file_size_mb = pdf_file.stat().st_size / (1024 * 1024)
        if file_size_mb > 100:
            self.logger.warning(f"Large PDF file detected: {file_size_mb:.1f}MB. "
                              "Conversion may take a while and use significant memory.")
        
        if not os.access(pdf_file, os.R_OK):
            raise PermissionError(f"No read permission for file: {pdf_path}")
        
        return pdf_file
    
    def validate_output_path(self, output_path: str) -> Path:
        """
        Validate and prepare the output path.
        
        Args:
            output_path: Desired output path for the PowerPoint file
            
        Returns:
            Path object for the output file
            
        Raises:
            PermissionError: If cannot write to the directory
            ValueError: If path is invalid
        """
        output_file = Path(output_path)
        
        if output_file.suffix.lower() != '.pptx':
            output_file = output_file.with_suffix('.pptx')
        
        parent_dir = output_file.parent
        if not parent_dir.exists():
            try:
                parent_dir.mkdir(parents=True, exist_ok=True)
                self.logger.info(f"Created output directory: {parent_dir}")
            except OSError as e:
                raise PermissionError(f"Cannot create output directory {parent_dir}: {e}")
        
        if not os.access(parent_dir, os.W_OK):
            raise PermissionError(f"No write permission for directory: {parent_dir}")
        
        if output_file.exists():
            self.logger.warning(f"Output file already exists and will be overwritten: {output_file}")
        
        return output_file
    
    def convert_pdf_to_images(self, pdf_path: Path) -> list:
        """
        Convert PDF to list of images with error handling.
        
        Args:
            pdf_path: Path to the PDF file
            
        Returns:
            List of PIL Image objects
            
        Raises:
            RuntimeError: If PDF conversion fails
        """
        try:
            self.logger.info(f"Converting PDF to images (DPI: {self.dpi}, "
                           f"Format: {self.image_format}, Threads: {self.thread_count})")
            
            images = convert_from_path(
                str(pdf_path), 
                dpi=self.dpi, 
                fmt=self.image_format,
                thread_count=self.thread_count
            )
            
            if not images:
                raise RuntimeError("PDF conversion resulted in no images")
            
            self.logger.info(f"Successfully converted PDF to {len(images)} images")
            return images
            
        except Exception as e:
            raise RuntimeError(f"Failed to convert PDF to images: {e}")
    
    def create_powerpoint(self, images: list, output_path: Path) -> None:
        """
        Create PowerPoint presentation from images.
        
        Args:
            images: List of PIL Image objects
            output_path: Path for the output PowerPoint file
            
        Raises:
            RuntimeError: If PowerPoint creation fails
        """
        try:
            prs = Presentation()
            blank_slide_layout = prs.slide_layouts[6]  # Blank slide layout
            
            for i, slideimg in enumerate(images):
                if i % self.progress_interval == 0:
                    self.logger.info(f"Processing slide {i + 1}/{len(images)}")
                
                try:
                    with BytesIO() as imagefile:
                        slideimg.save(imagefile, format=self.image_format.upper())
                        imagefile.seek(0)
                        width, height = slideimg.size
                        
                        # Set slide dimensions (only on first iteration)
                        if i == 0:
                            prs.slide_height = height * 9525  # Convert pixels to EMUs
                            prs.slide_width = width * 9525
                            self.logger.info(f"Set slide dimensions: {width}x{height} pixels")
                        
                        # Add slide
                        slide = prs.slides.add_slide(blank_slide_layout)
                        pic = slide.shapes.add_picture(
                            imagefile, 0, 0, 
                            width=width * 9525, 
                            height=height * 9525
                        )
                        
                except Exception as e:
                    self.logger.error(f"Failed to process slide {i + 1}: {e}")
                    raise RuntimeError(f"Failed to process slide {i + 1}: {e}")
            
            self.logger.info(f"Saving PowerPoint file: {output_path}")
            prs.save(str(output_path))
            self.logger.info("PowerPoint file saved successfully")
            
        except Exception as e:
            raise RuntimeError(f"Failed to create PowerPoint presentation: {e}")
    
    def convert(self, pdf_path: str, output_path: Optional[str] = None) -> Path:
        """
        Convert PDF to PowerPoint with comprehensive error handling.
        
        Args:
            pdf_path: Path to the input PDF file
            output_path: Optional path for output file (auto-generated if None)
            
        Returns:
            Path to the created PowerPoint file
            
        Raises:
            Various exceptions for different failure modes
        """
        pdf_file = self.validate_input_file(pdf_path)
        
        if output_path is None:
            output_path = str(pdf_file.with_suffix('.pptx'))
        
        output_file = self.validate_output_path(str(output_path))
        
        self.logger.info(f"Starting conversion: {pdf_file} -> {output_file}")
        
        try:
            images = self.convert_pdf_to_images(pdf_file)
            
            self.create_powerpoint(images, output_file)
            
            if not output_file.exists():
                raise RuntimeError("Output file was not created")
            
            file_size_mb = output_file.stat().st_size / (1024 * 1024)
            self.logger.info(f"Conversion completed successfully! "
                           f"Output: {output_file} ({file_size_mb:.1f}MB)")
            
            return output_file
            
        except Exception as e:
            self.logger.error(f"Conversion failed: {e}")
            if output_file.exists():
                try:
                    output_file.unlink()
                    self.logger.info("Cleaned up partial output file")
                except OSError:
                    pass
            raise


def setup_logging(verbose: bool = False) -> None:
    """Setup logging configuration."""
    level = logging.DEBUG if verbose else logging.INFO
    logging.basicConfig(
        level=level,
        format='%(asctime)s - %(levelname)s - %(message)s',
        datefmt='%Y-%m-%d %H:%M:%S'
    )


def main():
    """Main entry point with argument parsing."""
    parser = argparse.ArgumentParser(
        description='Convert PDF files to PowerPoint presentations',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  %(prog)s document.pdf                    # Convert to document.pptx
  %(prog)s document.pdf -o slides.pptx    # Convert with custom output name
  %(prog)s document.pdf --dpi 600         # High resolution conversion
  %(prog)s document.pdf --format jpeg     # Use JPEG compression
  %(prog)s document.pdf -v                # Verbose logging
        """
    )
    
    parser.add_argument('pdf_file', help='Input PDF file to convert')
    parser.add_argument('-o', '--output', help='Output PowerPoint file (default: same name as PDF)')
    parser.add_argument('--dpi', type=int, default=300, 
                       help='Resolution for PDF to image conversion (default: 300)')
    parser.add_argument('--format', choices=['png', 'jpeg', 'tiff'], default='png',
                       help='Image format for slides (default: png)')
    parser.add_argument('--threads', type=int, default=2,
                       help='Number of threads for PDF conversion (default: 2)')
    parser.add_argument('--progress', type=int, default=10,
                       help='Progress update interval in slides (default: 10)')
    parser.add_argument('-v', '--verbose', action='store_true',
                       help='Enable verbose logging')
    
    args = parser.parse_args()
    
    setup_logging(args.verbose)
    logger = logging.getLogger(__name__)
    
    try:
        converter = PDFToPowerPointConverter(
            dpi=args.dpi,
            image_format=args.format,
            thread_count=args.threads,
            progress_interval=args.progress
        )
        
        output_file = converter.convert(args.pdf_file, args.output)
        
        print(f"\n✅ Conversion completed successfully!")
        print(f"📄 Input:  {args.pdf_file}")
        print(f"📊 Output: {output_file}")
        
        return 0
        
    except KeyboardInterrupt:
        logger.error("Conversion interrupted by user")
        return 130
    except Exception as e:
        logger.error(f"Conversion failed: {e}")
        return 1


if __name__ == '__main__':
    sys.exit(main())
