#!/usr/bin/env python3
"""
PDF to PowerPoint Converter - Improved Version

This script converts PDF files into PowerPoint presentations by converting each page
to an image and embedding it in a slide. This ensures perfect visual fidelity without
font or formatting issues.

Author: Improved by Codegen
License: MIT
"""

import os
import sys
import argparse
import logging
from pathlib import Path
from typing import Optional, List

try:
    from PIL import Image
    from pdf2image import convert_from_path
    from pptx import Presentation
    from pptx.util import Inches
    from tqdm import tqdm
except ImportError as e:
    print(f"Error: Missing required dependency: {e}")
    print("Please install requirements with: pip install -r requirements.txt")
    sys.exit(1)


class PDFToPowerPointConverter:
    """Main converter class for PDF to PowerPoint conversion."""
    
    def __init__(self, dpi: int = 300, image_format: str = 'PNG', thread_count: int = 2):
        """
        Initialize the converter with configuration options.
        
        Args:
            dpi: Resolution for PDF to image conversion (default: 300)
            image_format: Image format for conversion (default: PNG)
            thread_count: Number of threads for PDF processing (default: 2)
        """
        self.dpi = dpi
        self.image_format = image_format.upper()
        self.thread_count = thread_count
        self.logger = self._setup_logger()
    
    def _setup_logger(self) -> logging.Logger:
        """Set up logging configuration."""
        logger = logging.getLogger(__name__)
        logger.setLevel(logging.INFO)
        
        if not logger.handlers:
            handler = logging.StreamHandler()
            formatter = logging.Formatter('%(levelname)s: %(message)s')
            handler.setFormatter(formatter)
            logger.addHandler(handler)
        
        return logger
    
    def validate_input_file(self, pdf_path: str) -> Path:
        """
        Validate the input PDF file.
        
        Args:
            pdf_path: Path to the PDF file
            
        Returns:
            Path object for the validated file
            
        Raises:
            FileNotFoundError: If file doesn't exist
            ValueError: If file is not a PDF
        """
        file_path = Path(pdf_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"PDF file not found: {pdf_path}")
        
        if not file_path.is_file():
            raise ValueError(f"Path is not a file: {pdf_path}")
        
        if file_path.suffix.lower() != '.pdf':
            raise ValueError(f"File is not a PDF: {pdf_path}")
        
        return file_path
    
    def get_output_path(self, input_path: Path, output_path: Optional[str] = None) -> Path:
        """
        Determine the output path for the PowerPoint file.
        
        Args:
            input_path: Path to the input PDF file
            output_path: Optional custom output path
            
        Returns:
            Path object for the output file
        """
        if output_path:
            output_file = Path(output_path)
            if output_file.suffix.lower() != '.pptx':
                output_file = output_file.with_suffix('.pptx')
        else:
            output_file = input_path.with_suffix('.pptx')
        
        return output_file
    
    def check_output_overwrite(self, output_path: Path, force: bool = False) -> bool:
        """
        Check if output file exists and handle overwrite confirmation.
        
        Args:
            output_path: Path to the output file
            force: Skip confirmation if True
            
        Returns:
            True if should proceed, False otherwise
        """
        if not output_path.exists():
            return True
        
        if force:
            self.logger.warning(f"Overwriting existing file: {output_path}")
            return True
        
        response = input(f"Output file '{output_path}' already exists. Overwrite? (y/N): ")
        return response.lower() in ['y', 'yes']
    
    def convert_pdf_to_images(self, pdf_path: Path) -> List[Image.Image]:
        """
        Convert PDF pages to images.
        
        Args:
            pdf_path: Path to the PDF file
            
        Returns:
            List of PIL Image objects
            
        Raises:
            Exception: If PDF conversion fails
        """
        try:
            self.logger.info(f"Converting PDF to images (DPI: {self.dpi})...")
            
            # Convert PDF to images
            images = convert_from_path(
                str(pdf_path),
                dpi=self.dpi,
                fmt='ppm',  # Internal format for processing
                thread_count=self.thread_count
            )
            
            self.logger.info(f"Successfully converted {len(images)} pages")
            return images
            
        except Exception as e:
            raise Exception(f"Failed to convert PDF to images: {str(e)}")
    
    def create_powerpoint(self, images: List[Image.Image], output_path: Path) -> None:
        """
        Create PowerPoint presentation from images.
        
        Args:
            images: List of PIL Image objects
            output_path: Path for the output PowerPoint file
            
        Raises:
            Exception: If PowerPoint creation fails
        """
        try:
            self.logger.info("Creating PowerPoint presentation...")
            
            # Create presentation
            prs = Presentation()
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            
            # Process each image with progress bar
            for i, image in enumerate(tqdm(images, desc="Processing slides", unit="slide")):
                # Convert image to bytes
                from io import BytesIO
                image_buffer = BytesIO()
                
                # Save image in specified format
                if self.image_format == 'PNG':
                    image.save(image_buffer, format='PNG', optimize=True)
                elif self.image_format == 'JPEG':
                    # Convert RGBA to RGB for JPEG
                    if image.mode == 'RGBA':
                        rgb_image = Image.new('RGB', image.size, (255, 255, 255))
                        rgb_image.paste(image, mask=image.split()[-1])
                        image = rgb_image
                    image.save(image_buffer, format='JPEG', quality=95, optimize=True)
                else:
                    image.save(image_buffer, format=self.image_format)
                
                image_buffer.seek(0)
                width, height = image.size
                
                # Set slide dimensions (convert pixels to EMUs: 1 pixel = 9525 EMUs at 96 DPI)
                # Adjust for actual DPI
                emu_per_pixel = 9525 * (96 / self.dpi)
                prs.slide_height = int(height * emu_per_pixel)
                prs.slide_width = int(width * emu_per_pixel)
                
                # Add slide
                slide = prs.slides.add_slide(blank_slide_layout)
                slide.shapes.add_picture(
                    image_buffer, 
                    0, 0, 
                    width=int(width * emu_per_pixel), 
                    height=int(height * emu_per_pixel)
                )
            
            # Save PowerPoint file
            self.logger.info(f"Saving PowerPoint file: {output_path}")
            prs.save(str(output_path))
            self.logger.info("Conversion completed successfully!")
            
        except Exception as e:
            raise Exception(f"Failed to create PowerPoint presentation: {str(e)}")
    
    def convert(self, pdf_path: str, output_path: Optional[str] = None, force: bool = False) -> None:
        """
        Main conversion method.
        
        Args:
            pdf_path: Path to the input PDF file
            output_path: Optional path for the output PowerPoint file
            force: Skip overwrite confirmation if True
        """
        try:
            # Validate input
            input_file = self.validate_input_file(pdf_path)
            output_file = self.get_output_path(input_file, output_path)
            
            # Check for overwrite
            if not self.check_output_overwrite(output_file, force):
                self.logger.info("Conversion cancelled by user")
                return
            
            # Convert PDF to images
            images = self.convert_pdf_to_images(input_file)
            
            # Create PowerPoint presentation
            self.create_powerpoint(images, output_file)
            
            # Display results
            file_size = output_file.stat().st_size / (1024 * 1024)  # MB
            self.logger.info(f"Output file: {output_file} ({file_size:.1f} MB)")
            
        except Exception as e:
            self.logger.error(f"Conversion failed: {str(e)}")
            sys.exit(1)


def main():
    """Main entry point for the script."""
    parser = argparse.ArgumentParser(
        description="Convert PDF files to PowerPoint presentations",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  %(prog)s document.pdf                    # Convert to document.pptx
  %(prog)s document.pdf -o presentation.pptx  # Specify output filename
  %(prog)s document.pdf --dpi 600 --force     # High resolution, overwrite existing
  %(prog)s document.pdf --format JPEG         # Use JPEG compression for smaller files
        """
    )
    
    parser.add_argument(
        'pdf_file',
        help='Path to the PDF file to convert'
    )
    
    parser.add_argument(
        '-o', '--output',
        help='Output PowerPoint file path (default: same name as PDF with .pptx extension)'
    )
    
    parser.add_argument(
        '--dpi',
        type=int,
        default=300,
        help='Resolution for PDF to image conversion (default: 300)'
    )
    
    parser.add_argument(
        '--format',
        choices=['PNG', 'JPEG', 'TIFF'],
        default='PNG',
        help='Image format for slides (default: PNG)'
    )
    
    parser.add_argument(
        '--threads',
        type=int,
        default=2,
        help='Number of threads for PDF processing (default: 2)'
    )
    
    parser.add_argument(
        '--force',
        action='store_true',
        help='Overwrite output file without confirmation'
    )
    
    parser.add_argument(
        '--verbose',
        action='store_true',
        help='Enable verbose logging'
    )
    
    parser.add_argument(
        '--version',
        action='version',
        version='PDF to PowerPoint Converter v2.0'
    )
    
    args = parser.parse_args()
    
    # Set up logging level
    if args.verbose:
        logging.getLogger().setLevel(logging.DEBUG)
    
    # Create converter and run conversion
    converter = PDFToPowerPointConverter(
        dpi=args.dpi,
        image_format=args.format,
        thread_count=args.threads
    )
    
    converter.convert(
        pdf_path=args.pdf_file,
        output_path=args.output,
        force=args.force
    )


if __name__ == '__main__':
    main()
