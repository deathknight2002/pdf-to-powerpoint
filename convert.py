import argparse
import os
import sys
import logging
from PIL import Image
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO

EMU_PER_PIXEL = 9525

def parse_arguments():
    parser = argparse.ArgumentParser(
        description='Convert PDF files to PowerPoint presentations using images'
    )
    parser.add_argument(
        'pdf_file',
        help='Path to the PDF file to convert'
    )
    parser.add_argument(
        '--dpi',
        type=int,
        default=300,
        help='DPI for image conversion (default: 300)'
    )
    parser.add_argument(
        '--format',
        default='png',
        choices=['png', 'tiff', 'jpeg', 'ppm'],
        help='Image format for conversion (default: png)'
    )
    parser.add_argument(
        '--threads',
        type=int,
        default=2,
        help='Number of threads for processing (default: 2)'
    )
    parser.add_argument(
        '--batch-size',
        type=int,
        default=10,
        help='Number of pages to process in each batch (default: 10)'
    )
    return parser.parse_args()

def validate_input(pdf_file):
    if not os.path.exists(pdf_file):
        raise FileNotFoundError(f"PDF file not found: {pdf_file}")
    
    if not os.path.isfile(pdf_file):
        raise ValueError(f"Path is not a file: {pdf_file}")
    
    if not pdf_file.lower().endswith('.pdf'):
        raise ValueError(f"File is not a PDF: {pdf_file}")
    
    if not os.access(pdf_file, os.R_OK):
        raise PermissionError(f"Cannot read PDF file: {pdf_file}")

def process_pdf_in_batches(pdf_file, dpi, img_format, thread_count, batch_size):
    logging.info("Starting PDF conversion...")
    
    try:
        slideimgs = convert_from_path(
            pdf_file, 
            dpi=dpi, 
            fmt='ppm', 
            thread_count=thread_count
        )
    except Exception as e:
        raise RuntimeError(f"Failed to convert PDF to images: {e}")
    
    total_slides = len(slideimgs)
    logging.info(f"PDF converted to {total_slides} images")
    
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    
    for batch_start in range(0, total_slides, batch_size):
        batch_end = min(batch_start + batch_size, total_slides)
        logging.info(f"Processing slides {batch_start + 1}-{batch_end} of {total_slides}")
        
        for i in range(batch_start, batch_end):
            slideimg = slideimgs[i]
            
            try:
                with BytesIO() as imagefile:
                    slideimg.save(imagefile, format=img_format)
                    imagefile.seek(0)
                    width, height = slideimg.size
                    
                    prs.slide_height = height * EMU_PER_PIXEL
                    prs.slide_width = width * EMU_PER_PIXEL
                    
                    slide = prs.slides.add_slide(blank_slide_layout)
                    slide.shapes.add_picture(
                        imagefile, 
                        0, 
                        0, 
                        width=width * EMU_PER_PIXEL, 
                        height=height * EMU_PER_PIXEL
                    )
                    
            except Exception as e:
                raise RuntimeError(f"Failed to process slide {i + 1}: {e}")
        
        progress = (batch_end / total_slides) * 100
        logging.info(f"Progress: {progress:.1f}% complete")
    
    return prs

def main():
    logging.basicConfig(
        level=logging.INFO,
        format='%(message)s'
    )
    
    try:
        args = parse_arguments()
        
        validate_input(args.pdf_file)
        
        logging.info(f"Converting file: {args.pdf_file}")
        logging.info(f"Settings: DPI={args.dpi}, Format={args.format}, Threads={args.threads}, Batch Size={args.batch_size}")
        
        prs = process_pdf_in_batches(
            args.pdf_file,
            args.dpi,
            args.format,
            args.threads,
            args.batch_size
        )
        
        base_name = args.pdf_file.rsplit('.pdf', 1)[0]
        output_file = base_name + '.pptx'
        
        logging.info(f"Saving file: {output_file}")
        prs.save(output_file)
        logging.info("Conversion complete! :)")
        
    except FileNotFoundError as e:
        logging.error(f"File error: {e}")
        sys.exit(1)
    except ValueError as e:
        logging.error(f"Input error: {e}")
        sys.exit(1)
    except PermissionError as e:
        logging.error(f"Permission error: {e}")
        sys.exit(1)
    except RuntimeError as e:
        logging.error(f"Conversion error: {e}")
        sys.exit(1)
    except Exception as e:
        logging.error(f"Unexpected error: {e}")
        sys.exit(1)

if __name__ == "__main__":
    main()
